"""Build YouTube Citations PowerPoint from slide data."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import os

# Samsung brand colors
SAMSUNG_BLUE = RGBColor(0x14, 0x28, 0xA0)
SAMSUNG_DARK_BLUE = RGBColor(0x0D, 0x1A, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
BG_SECTION = RGBColor(0xFA, 0xFA, 0xFA)
GREEN = RGBColor(0x96, 0xD5, 0x51)
RED = RGBColor(0xFF, 0x44, 0x38)
NEUTRAL_GRAY = RGBColor(0x9E, 0x9E, 0x9E)
YT_RED = RGBColor(0xFF, 0x00, 0x00)
INSIGHT_GREEN = RGBColor(0x2E, 0x7D, 0x32)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_DISPLAY = "Samsung Sharp Sans"
FONT_BODY = "Samsung One"
FONT_FALLBACK = "Calibri"


def add_title_bar(slide, title, subtitle, badge):
    """Add Samsung blue header bar to a slide."""
    # Blue header background
    shp = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.95)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = SAMSUNG_BLUE
    shp.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_DISPLAY
    if subtitle:
        run2 = p.add_run()
        run2.text = f"  {subtitle}"
        run2.font.size = Pt(14)
        run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)
        run2.font.name = FONT_BODY

    # Badge
    badge_box = slide.shapes.add_textbox(Inches(10), Inches(0.25), Inches(3), Inches(0.45))
    tf2 = badge_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run3 = p2.add_run()
    run3.text = badge
    run3.font.size = Pt(12)
    run3.font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)
    run3.font.name = FONT_BODY


def add_kpi_card(slide, left, top, width, value, label, subtext):
    """Add a KPI card shape."""
    height = Inches(1.35)
    # Card background
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_SECTION
    shp.line.fill.background()

    # Blue top accent
    accent = slide.shapes.add_shape(1, left, top, width, Pt(4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = SAMSUNG_BLUE
    accent.line.fill.background()

    # Value
    val_box = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.55))
    tf = val_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = SAMSUNG_BLUE
    run.font.name = FONT_DISPLAY

    # Label
    lbl_box = slide.shapes.add_textbox(left, top + Inches(0.7), width, Inches(0.3))
    tf2 = lbl_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label.upper()
    run2.font.size = Pt(11)
    run2.font.bold = True
    run2.font.color.rgb = BLACK
    run2.font.name = FONT_DISPLAY

    # Subtext
    sub_box = slide.shapes.add_textbox(left, top + Inches(1.0), width, Inches(0.25))
    tf3 = sub_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = subtext
    run3.font.size = Pt(9)
    run3.font.color.rgb = GRAY
    run3.font.name = FONT_BODY


def add_insight_card(slide, left, top, width, title_text, body_text, height=Inches(1.1)):
    """Add an insight card with blue left border."""
    # Background
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.fill.background()

    # Blue left accent
    accent = slide.shapes.add_shape(1, left, top, Pt(4), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = SAMSUNG_BLUE
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.25))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = SAMSUNG_BLUE
    run.font.name = FONT_DISPLAY

    # Body
    body_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.35), width - Inches(0.3), height - Inches(0.4))
    tf2 = body_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = body_text
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY
    run2.font.name = FONT_BODY


def add_horizontal_bar_chart(slide, left, top, width, height, data, title_text, color=SAMSUNG_BLUE):
    """Add a horizontal bar chart."""
    chart_data = CategoryChartData()
    chart_data.categories = [d[0] for d in data]
    chart_data.add_series("Citations", [d[1] for d in data])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = title_text
    for run in chart.chart_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.name = FONT_DISPLAY

    plot = chart.plots[0]
    plot.gap_width = 80

    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color if isinstance(color, RGBColor) else SAMSUNG_BLUE
    series.data_labels.show_value = True
    series.data_labels.font.size = Pt(9)
    series.data_labels.font.name = FONT_BODY
    series.data_labels.number_format = '#,##0'

    # Style axes
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.name = FONT_BODY
    cat_axis.has_major_gridlines = False

    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.has_major_gridlines = False

    return chart_frame


def add_table(slide, left, top, width, height, headers, rows):
    """Add a styled table."""
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height
    )
    table = table_shape.table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = SAMSUNG_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_BODY
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        bg = WHITE if row_idx % 2 == 0 else BG_SECTION
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.name = FONT_BODY
                paragraph.font.color.rgb = BLACK
                if col_idx == 0:
                    paragraph.alignment = PP_ALIGN.CENTER
                elif headers[col_idx] in ("Citations", "Views", "Likes", "Videos"):
                    paragraph.alignment = PP_ALIGN.RIGHT
                    if headers[col_idx] == "Citations":
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = SAMSUNG_BLUE

    return table_shape


def add_timestamp_card(slide, left, top, width, time_str, title, channel, citations, quote):
    """Add a timestamp citation card."""
    height = Inches(1.4)

    # Card background
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.fill.background()

    # Red left accent
    accent = slide.shapes.add_shape(1, left, top, Pt(4), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = YT_RED
    accent.line.fill.background()

    # Timestamp badge
    badge_w = Inches(0.6)
    badge_h = Inches(0.3)
    badge = slide.shapes.add_shape(1, left + Inches(0.12), top + Inches(0.12), badge_w, badge_h)
    badge.fill.solid()
    badge.fill.fore_color.rgb = YT_RED
    badge.line.fill.background()
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = badge.text_frame.paragraphs[0].add_run()
    run.text = time_str
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_DISPLAY

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.85), top + Inches(0.08), width - Inches(1.0), Inches(0.25))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title[:55] + ("..." if len(title) > 55 else "")
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = BLACK
    r.font.name = FONT_BODY

    # Channel + citations
    meta_box = slide.shapes.add_textbox(left + Inches(0.85), top + Inches(0.32), width - Inches(1.0), Inches(0.2))
    tf2 = meta_box.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = f"{channel}  ·  {citations} citations"
    r2.font.size = Pt(8)
    r2.font.color.rgb = GRAY
    r2.font.name = FONT_BODY

    # Quote
    quote_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.58), width - Inches(0.35), height - Inches(0.65))
    tf3 = quote_box.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = f'"{quote[:140]}{"..." if len(quote) > 140 else ""}"'
    r3.font.size = Pt(8)
    r3.font.italic = True
    r3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    r3.font.name = FONT_BODY


def add_quote_card(slide, left, top, width, quote_text, source, video, tag):
    """Add a quote card."""
    height = Inches(1.45)

    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.fill.background()

    accent = slide.shapes.add_shape(1, left, top, Pt(4), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = SAMSUNG_BLUE
    accent.line.fill.background()

    # Quote text
    q_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.75))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f'"{quote_text}"'
    r.font.size = Pt(9)
    r.font.italic = True
    r.font.color.rgb = BLACK
    r.font.name = FONT_BODY

    # Source
    s_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.9), width - Inches(1.0), Inches(0.2))
    tf2 = s_box.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = f"— {source}, {video}"
    r2.font.size = Pt(8)
    r2.font.color.rgb = GRAY
    r2.font.name = FONT_BODY

    # Tag badge
    tag_box = slide.shapes.add_shape(1, left + width - Inches(0.9), top + height - Inches(0.35), Inches(0.8), Inches(0.25))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = SAMSUNG_BLUE
    tag_box.line.fill.background()
    tag_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tr = tag_box.text_frame.paragraphs[0].add_run()
    tr.text = tag
    tr.font.size = Pt(8)
    tr.font.bold = True
    tr.font.color.rgb = WHITE
    tr.font.name = FONT_BODY


def build_slide_1(prs):
    """YouTube Citations Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_title_bar(slide, "YouTube Citations Overview", "AI Model Video References", "YouTube Citations | Feb 2026")

    # KPI Cards
    kpis = [
        ("164", "Videos Analyzed", "Unique YouTube videos cited"),
        ("86", "Unique Channels", "Distinct YouTube creators"),
        ("3,175", "Total Citations", "YouTube links in AI responses"),
        ("4,469", "Comments Collected", "User comments analyzed"),
        ("18.2M", "Total Views", "Aggregate video view count"),
    ]
    card_w = Inches(2.35)
    gap = Inches(0.18)
    start_x = Inches(0.5)
    for i, (val, lbl, sub) in enumerate(kpis):
        add_kpi_card(slide, start_x + i * (card_w + gap), Inches(1.15), card_w, val, lbl, sub)

    # Bar chart
    prompt_data = [
        ("Mini-LED vs QLED", 66),
        ("Mini-LED tv explained", 65),
        ("best-reviewed Mini-LED tvs", 52),
        ("Mini-LED vs OLED", 52),
        ("Mini-LED vs Micro RGB differences", 48),
        ("is OLED worth extra over Mini-LED", 46),
        ("is Mini-LED better than OLED", 46),
        ("what is Mini-LED", 43),
        ("is Mini-LED worth it", 42),
        ("Mini-LED gaming performance", 40),
    ]
    add_horizontal_bar_chart(
        slide, Inches(0.5), Inches(2.75), Inches(8.0), Inches(4.3),
        prompt_data, "Top Mini-LED Prompts Triggering YouTube Citations"
    )

    # Insight cards
    insights = [
        ("Comparison Queries Dominate", '"Mini-LED vs QLED" and "Mini-LED vs OLED" are the top prompts. AI models cite YouTube videos when users ask head-to-head technology questions.'),
        ("Educational Content Cited Heavily", '"What is Mini-LED" and "Mini-LED tv explained" drive 108 combined citations — AI models rely on YouTube explainers to educate users.'),
        ("Purchase Intent Signals", '"Is Mini-LED worth it", "best-reviewed Mini-LED tvs", and gaming performance queries show AI models surfacing video content at the decision stage.'),
    ]
    for i, (title, body) in enumerate(insights):
        add_insight_card(slide, Inches(8.8), Inches(2.75) + i * Inches(1.5), Inches(4.0), title, body, Inches(1.35))


def build_slide_2(prs):
    """Influencer Channel Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Influencer Channel Analysis", "Top YouTube Channels Cited for Mini-LED Content", "YouTube Citations | Feb 2026")

    # Table
    headers = ["#", "Channel", "Subs", "Videos", "Citations", "Views"]
    rows = [
        [1, "RTINGS Home Theater", "398K", 7, 111, "719K"],
        [2, "Digital Trends", "1.37M", 7, 105, "3.8M"],
        [3, "Explaining Tech Like You're Five", "78.5K", 2, 102, "901K"],
        [4, "Switch and Click", "1.09M", 1, 80, "417K"],
        [5, "ForeMost Picks", "2.9K", 1, 69, "10.6K"],
        [6, "The Viewing Angle", "572K", 4, 64, "1.5M"],
        [7, "Jon Rettinger", "1.65M", 7, 47, "3.3M"],
        [8, "The Display Guy", "50.8K", 7, 42, "612K"],
        [9, "Toasty Bros", "893K", 1, 29, "118K"],
        [10, "all things tech", "279K", 1, 21, "435K"],
        [11, "NextGen Picks", "2.3K", 3, 19, "2.3K"],
        [12, "New Gen Home", "2.8K", 3, 18, "33.3K"],
        [13, "Blumoo", "11.8K", 1, 18, "46"],
        [14, "Smart Home Sounds", "154K", 5, 16, "460K"],
        [15, "B The Installer", "276K", 2, 15, "125K"],
    ]
    add_table(slide, Inches(0.4), Inches(1.15), Inches(7.5), Inches(4.8), headers, rows)

    # Bar chart (top 10)
    chart_data = [(r[1][:20], r[4]) for r in rows[:10]]
    add_horizontal_bar_chart(
        slide, Inches(8.2), Inches(1.15), Inches(4.8), Inches(3.2),
        chart_data, "Citations by Channel"
    )

    # Insights
    insights = [
        ("RTINGS & Digital Trends Lead", "The two largest tech review outlets combine for 216 Mini-LED citations across 14 videos."),
        ("Small Channels Punch Above Weight", "ForeMost Picks (2.9K subs) earns 69 citations from a single video, outperforming Jon Rettinger (1.65M subs, 47 citations)."),
        ("Explainer Channels Dominate", "Explaining Tech Like You're Five ranks #3 with 102 citations from just 2 videos. AI favors accessible Mini-LED explainers."),
    ]
    for i, (title, body) in enumerate(insights):
        add_insight_card(slide, Inches(8.2), Inches(4.55) + i * Inches(1.05), Inches(4.8), title, body, Inches(0.95))


def build_slide_3(prs):
    """Most Cited Videos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Most Cited Videos", "Top Mini-LED YouTube Content Referenced by AI", "YouTube Citations | Feb 2026")

    headers = ["#", "Title", "Channel", "Citations", "Views", "Published"]
    rows = [
        [1, "The Best TVs To Buy in 2025 | OLED, QLED, Mini LED", "RTINGS", 93, "167K", "Oct 2025"],
        [2, "OLED, QLED, or Mini-LED? Which TV to Buy in 2025", "Digital Trends", 91, "2.87M", "Apr 2025"],
        [3, "Mini LED vs OLED Monitors (don't waste your money)", "Switch & Click", 80, "417K", "Nov 2025"],
        [4, "OLED vs QLED vs Mini LED in 2026 | Best Choice!", "ForeMost Picks", 69, "10.6K", "Jan 2026"],
        [5, "OLED vs QLED: Don't Get Tricked in the Store", "Explain Tech x5", 63, "834K", "Nov 2025"],
        [6, "Mini-LED vs OLED — Which One Should YOU Buy?", "Explain Tech x5", 39, "66.9K", "Dec 2025"],
        [7, "The new king for 2026! QLED vs Mini-LED TVs", "Viewing Angle", 39, "402K", "Oct 2025"],
        [8, "OLED Vs Mini LED - Pros & Cons", "Toasty Bros", 29, "118K", "Nov 2025"],
        [9, "Did you buy the wrong TV? OLED vs Mini LED", "Jon Rettinger", 24, "428K", "Nov 2025"],
        [10, "Mini LED vs OLED TVs 2026! Don't make a mistake", "Viewing Angle", 23, "488K", "Jul 2025"],
        [11, "Best TVs of 2025 (OLED & Mini LED)", "Display Guy", 22, "8.4K", "Dec 2025"],
        [12, "Mini LED vs Standard LED TVs — BIG Difference", "all things tech", 21, "435K", "Nov 2023"],
        [13, "Best Mini LED TVs 2026 - Buying Guide", "Blumoo", 18, "46", "Jan 2026"],
        [14, "OLED vs QLED vs Mini-LED: Which TV in 2026?", "New Gen Home", 16, "2.7K", "Jan 2026"],
        [15, "Best Mini LED TVs (End of Year 2025)", "Royalty Consumer", 15, "61.7K", "Sep 2025"],
    ]
    add_table(slide, Inches(0.4), Inches(1.15), Inches(12.5), Inches(4.5), headers, rows)

    # Insights row
    insights = [
        ("Comparison Format Dominates", "11 of 15 top videos are direct comparisons (OLED vs Mini-LED, QLED vs Mini-LED). AI overwhelmingly favors head-to-head content."),
        ("Viral Reach + AI Citations", "Digital Trends' guide has 2.87M views AND 91 citations. High-production comparison content earns both audience reach and AI trust."),
        ("Evergreen Content Stays Cited", "all things tech's 2023 explainer still earns 21 citations in 2026. Well-structured Mini-LED educational content maintains AI relevance for years."),
    ]
    for i, (title, body) in enumerate(insights):
        add_insight_card(slide, Inches(0.4) + i * Inches(4.25), Inches(5.85), Inches(4.0), title, body, Inches(1.1))


def build_slide_4(prs):
    """Timestamp Deep Dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Timestamp Deep Dive", "AI Models Cite Specific Mini-LED Video Moments", "YouTube Citations | Feb 2026")

    # Intro banner
    banner = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.55))
    banner.fill.solid()
    banner.fill.fore_color.rgb = SAMSUNG_BLUE
    banner.line.fill.background()
    tf = banner.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "12 Mini-LED videos cited with specific timestamps — AI models pinpoint exact moments where Mini-LED technology is explained or compared"
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = FONT_DISPLAY

    # Timestamp cards (2 columns, 6 rows)
    cards = [
        ("0:13", "Mini-LED vs OLED — Which One Should YOU Buy?", "Explaining Tech Like You're Five", 39, "Most shoppers can't even tell which is which until someone points out the differences."),
        ("0:36", "Did you buy the wrong TV? OLED vs Mini LED", "Jon Rettinger", 24, "Movie nights, Rams games with every light on, gaming sessions that have gone way too late."),
        ("7:27", "Best TVs of 2025 (OLED & Mini LED)", "The Display Guy", 22, "The best mini-LED TV of 2025 has got to go to the Sony Bravia 9."),
        ("1:07", "Mini LED vs Standard LED TVs — BIG Difference", "all things tech", 21, "Then you get mini LEDs where there's thousands of LEDs giving you better brightness."),
        ("6:49", "My pick! The best type of TV in 2026!", "Jon Rettinger", 12, "Mini-LED sits in the sweet spot. Thousands of tiny LEDs controlling the backlight precisely."),
        ("0:36", "RGB and SQD-Mini LED — What's new in 2026?", "Tech Steve", 8, "A traditional mini-LED TV uses a blue backlight with quantum dots to create red and green."),
        ("0:27", "Best TVs of 2026 EXPOSED: OLED vs Mini-LED!", "NextGen Picks", 7, "Every single TV technology has a failure point."),
        ("0:08", "Mini-LED vs OLED in 2026 | Sony Bravia 9 Review", "Jeff Rauseo", 6, "Today, I want to talk about this brand new TV."),
        ("0:03", "Amazon Fire TV Omni Mini-LED TV review", "TechRadar", 5, "Amazon's most advanced own-brand TV. The brightest we've seen."),
        ("2:01", "What Type of TV? OLED, Mini-LED, RGB (2026)", "Tom's Guide", 4, "Samsung calls its mini-LEDs Neo QLED — marketing terms muddy the waters."),
        ("1:04", "Top 10 Best Mini-LED TVs in 2025", "Factic", 4, "Let's dive into the top 10 mini-LED TVs you need to know before buying."),
        ("0:20", "Can Hisense Knock Sony Off the Throne?", "Digital Trends", 3, "I've had the Bravia 9 for a few weeks, and I really want to compare these two."),
    ]

    col_w = Inches(6.0)
    for i, (time_str, title, channel, cit, quote) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * (col_w + Inches(0.3))
        y = Inches(1.9) + row * Inches(1.5)
        if y + Inches(1.4) > Inches(7.2):
            break
        add_timestamp_card(slide, x, y, col_w, time_str, title, channel, cit, quote)


def build_slide_5(prs):
    """Transcript Insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "What's Being Said", "Transcript Analysis from Top Mini-LED Cited Videos", "YouTube Citations | Feb 2026")

    # Themes chart
    theme_data = [
        ("Mini-LED vs OLED Comparisons", 52),
        ("Brightness & HDR Performance", 46),
        ("Dimming Zones & Local Dimming", 42),
        ("Price & Value Proposition", 38),
        ("Blooming & Halo Effects", 35),
        ("Gaming Performance", 31),
        ("Room Lighting Considerations", 28),
        ("Samsung Neo QLED Branding", 24),
        ("Contrast & Black Levels", 22),
        ("RGB Mini-LED & Next Gen", 18),
    ]
    add_horizontal_bar_chart(
        slide, Inches(0.4), Inches(1.15), Inches(6.2), Inches(3.5),
        theme_data, "Key Themes in Mini-LED Transcripts"
    )

    # Word frequency chart
    word_data = [
        ("Mini-LED", 892),
        ("OLED", 756),
        ("backlight", 634),
        ("dimming zones", 587),
        ("brightness", 523),
        ("Samsung", 498),
        ("contrast", 467),
        ("Neo QLED", 412),
        ("blooming", 389),
        ("HDR", 356),
        ("Sony Bravia 9", 334),
        ("TCL", 312),
        ("Hisense", 298),
        ("quantum dots", 276),
        ("QLED", 245),
    ]
    add_horizontal_bar_chart(
        slide, Inches(6.9), Inches(1.15), Inches(6.0), Inches(3.5),
        word_data, "Most Frequent Terms"
    )

    # Quote cards (2 columns, 3 rows)
    quotes = [
        ("Mini-LED sits in the sweet spot. Thousands of tiny LEDs controlling the backlight way more precisely.", "Jon Rettinger", "Best type of TV in 2026", "Mini-LED"),
        ("Most shoppers can't even tell which is which until someone points out the differences.", "Explaining Tech x5", "Mini-LED vs OLED", "Comparison"),
        ("I'm testing a $250 monitor that's just as good as a $1,500 OLED. Did I really waste $1,000?", "Switch and Click", "Mini LED vs OLED Monitors", "Value"),
        ("Samsung calls its mini-LEDs Neo QLED — this is where marketing terms start to muddy the waters.", "Tom's Guide", "OLED, Mini-LED, RGB Compared", "Branding"),
        ("You've seen QLED and Mini-LED plastered everywhere. My job is to help this make sense.", "The Viewing Angle", "QLED vs Mini-LED TVs", "Education"),
        ("A traditional mini-LED TV uses a blue backlight with quantum dots. This Hisense uses three colors on each pixel.", "Tech Steve", "RGB and SQD-Mini LED", "Technology"),
    ]

    # Section title
    qt = slide.shapes.add_textbox(Inches(0.4), Inches(4.8), Inches(6), Inches(0.35))
    tf = qt.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Representative Quotes from Top-Cited Mini-LED Transcripts"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.name = FONT_DISPLAY
    r.font.color.rgb = BLACK

    card_w = Inches(4.1)
    for i, (text, source, video, tag) in enumerate(quotes):
        col = i % 3
        row = i // 3
        x = Inches(0.4) + col * (card_w + Inches(0.2))
        y = Inches(5.25) + row * Inches(1.55)
        add_quote_card(slide, x, y, card_w, text, source, video, tag)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)

    out_dir = os.path.join(os.path.dirname(__file__), "..", "slides", "youtube_citations")
    out_path = os.path.join(out_dir, "YouTube_Citations_Mini_LED.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
